"""
Build the 2026 Field Day presentation. Run from anywhere:

    python3 site-setup/build-presentation.py

Output is written next to this script as 2026-field-day-presentation.pptx.
Edit the slide content below and re-run to regenerate.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
OUT = os.path.join(HERE, "2026-field-day-presentation.pptx")
LOGO_STRIP = os.path.join(REPO, "brochures", "build", "four_club_strip.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LAY_TITLE = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]
LAY_SECTION = prs.slide_layouts[5]
LAY_BLANK = prs.slide_layouts[6]

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
MUTED = RGBColor(0x55, 0x55, 0x55)


def add_title_slide(title, subtitle):
    s = prs.slides.add_slide(LAY_TITLE)
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    for para in s.shapes.title.text_frame.paragraphs:
        for r in para.runs:
            r.font.color.rgb = ACCENT
            r.font.bold = True
    return s


def add_bullets(title, bullets, subtitle=None):
    s = prs.slides.add_slide(LAY_CONTENT)
    s.shapes.title.text = title
    for r in s.shapes.title.text_frame.paragraphs[0].runs:
        r.font.color.rgb = ACCENT
        r.font.bold = True
    body = s.placeholders[1].text_frame
    body.clear()
    if subtitle:
        body.text = subtitle
        p0 = body.paragraphs[0]
        for r in p0.runs:
            r.font.italic = True
            r.font.color.rgb = MUTED
            r.font.size = Pt(16)
        first = False
    else:
        first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = text
        p.level = level
    return s


def add_section_divider(title, subtitle=None):
    s = prs.slides.add_slide(LAY_SECTION)
    s.shapes.title.text = title
    for r in s.shapes.title.text_frame.paragraphs[0].runs:
        r.font.color.rgb = ACCENT
        r.font.bold = True
        r.font.size = Pt(44)
    if subtitle:
        tb = s.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
        tf = tb.text_frame
        tf.text = subtitle
        for r in tf.paragraphs[0].runs:
            r.font.italic = True
            r.font.color.rgb = MUTED
            r.font.size = Pt(24)
    return s


# ---------------------------------------------------------------- slides
# 1 - Title
s = add_title_slide(
    "2026 ARRL Field Day",
    "UPARC · SPARC · WORMs · CARS  —  Contest Callsign W4TA  —  Class 3A",
)
if os.path.exists(LOGO_STRIP):
    s.shapes.add_picture(LOGO_STRIP, Inches(2.5), Inches(5.5), height=Inches(1.4))

# 2 - What is Field Day
add_bullets(
    "What is ARRL Field Day?",
    [
        "Ham radio's annual “open house” — a public demonstration of amateur radio.",
        "An exercise in setting up radio communications under less-than-ideal conditions: temporary antennas, generators, portable gear.",
        "Every June, 40,000+ hams across North America operate from public sites.",
        "Combines public service, emergency preparedness, community outreach, and technical skill.",
        "Annual event since 1933 — the most popular event in ham radio.",
    ],
)

# 3 - Why we do it
add_bullets(
    "Why we do it",
    [
        "Preparedness: prove we can stand up working stations from scratch when normal infrastructure isn’t available.",
        "Skill: practice operating in pile-ups, on every band and mode we have antennas for.",
        "Community: four clubs (UPARC, SPARC, WORMs, CARS) working a single event.",
        "Outreach: visitors welcome — Field Day is a great place to see ham radio in action.",
        "Fun: a 24-hour contest with friends, food, and a lot of QSOs.",
    ],
)

# 4 - When/where/schedule
add_bullets(
    "When, Where & Schedule",
    [
        ("Where: Clearwater Fire Station #48", 0),
        ("1700 N Belcher Rd., Clearwater, FL 33765", 1),
        ("When: June 27–28, 2026 (Saturday – Sunday)", 0),
        ("Setup begins: 8:00 AM Saturday, June 27 (lunch around noon)", 1),
        ("Operating begins: 2:00 PM Saturday, June 27", 1),
        ("Operating ends / teardown: 2:00 PM Sunday, June 28", 1),
        ("ARRL Class: 3A  (three simultaneous transmitters, portable/club)", 0),
        ("Operation is indoors (except the Satellite station).", 0),
        ("Dry-run: Saturday June 20, 9:00 AM, Clearwater Fire Training Center", 0),
        ("Meal signup: www.uparc.org — by June 22 at 11:59 PM", 0),
    ],
)

# 5 - Modes & Bands
add_bullets(
    "Operating Modes & Bands",
    [
        "Modes: CW, SSB, FT8 (digital)",
        "Bands: 6 m – 80 m",
        "Each station has a primary mode; other modes available if conditions dictate.",
        "Bring your own headset/mic gear (3.5 mm or 1/4 in plug). We do not share headphones or mic headsets for sanitary reasons.",
    ],
)

# 6 - Station Lineup divider
add_section_divider(
    "Station Lineup",
    "Three HF positions + one VHF position + Talk-In + Satellite",
)

# 7 - Station details
add_bullets(
    "Stations",
    [
        ("Station 1 — Elecraft K4 (SSB)", 0),
        ("K4’s CESSB feature gives ~6–8 dB of extra effective talk power. Single K4 used on SSB.", 1),
        ("Station 2 — Elecraft K3S (SSB / Digital incl. FT8)", 0),
        ("Adds an external VGA monitor for digital-mode waterfall and decoders.", 1),
        ("Station 3 — Elecraft K3 (CW)", 0),
        ("WinKey Mini + Bencher paddle (supplied by NY4I). Also runs SSB / Digital if needed.", 1),
        ("VHF Station — Yaesu FT-897 (6 m)", 0),
        ("Brought and operated by Paul KC4YDY.", 1),
        ("Talk-In Station — informational, not on the logging network.", 0),
    ],
)

# 8 - Network / Logging architecture
add_bullets(
    "Logging Network & Stats",
    [
        ("All HF stations run TR4W (Windows logger) on their laptops.", 0),
        ("Laptops join the TRLOG WiFi (10.0.0.0/24) via the Linksys router.", 0),
        ("Raspberry Pi 400 (wired Ethernet, 192.168.0.100) is the server:", 0),
        ("TR4W Server (Python) — multi-op duplicate checks across stations.", 1),
        ("n1mm_view — generates stats images (band/score/rate) on a RAM disk.", 1),
        ("Local web UI on :8080 — viewable on the on-site HDMI monitor.", 1),
        ("rsync publisher — pushes stats images to UPARC and SPARC websites.", 1),
        ("Pi 400 + Linksys router are both protected by a UPS.", 0),
        ("Replaces the old Windows TR4W Server box — single Pi handles everything now.", 0),
    ],
    subtitle="Tom NY4I",
)

# 9 - Antennas
add_bullets(
    "Antennas & Patch Panel",
    [
        ("All antenna feedlines come back to a portable patch panel in the radio room.", 0),
        ("Radios are connected to specific antennas through this patch panel.", 0),
        ("Do NOT change coax routing without approval from the assigned Band Boss.", 0),
        ("Station 3 (CW) uses a 65 ft RG-8X coax run around the room (east wall, top wall) to reach the patch panel.", 0),
    ],
    subtitle="Ryan AF4O",
)

# 10 - Satellite
add_bullets(
    "Satellite Station",
    [
        "ICOM IC-9700 transceiver",
        "Yaesu AZ/EL rotator",
        "Arrow yagi antenna",
        "Satellite control / tracking device",
        "Operated outdoors (only station not in the radio room).",
    ],
)

# 11 - Education
add_bullets(
    "Education & Training Sessions",
    [
        "Welcome to Field Day — for all visitors",
        "Ham Radio Basics",
        "Ham Radio Awards & Achievements",
        "How to build an EFHW kit antenna",
        "How to use an Antenna Analyzer",
        "All sessions led by Fred W2SUB.",
    ],
    subtitle="Fred W2SUB",
)

# 12 - Food
add_bullets(
    "Hospitality (a.k.a. Food)",
    [
        ("Saturday breakfast: donuts, coffee", 0),
        ("Saturday lunch: hamburgers, hotdogs, salad, chips", 0),
        ("Saturday dinner: TBD — bring a side dish to share", 0),
        ("Sunday breakfast: bagels (Tom), coffee", 0),
        ("Sunday lunch: leftovers", 0),
        ("Water, tea and sodas provided throughout the event.", 0),
        ("If you have special dietary needs, please plan to provide for yourself.", 0),
        ("Meal signup: www.uparc.org — by June 22 at 11:59 PM.", 0),
    ],
    subtitle="Ron W4RFA",
)

# 13 - Visitors / Signup
add_bullets(
    "Visitors & How to Get Involved",
    [
        "Visitors are very welcome — Field Day is a public outreach event.",
        "Club members: meal signup at www.uparc.org no later than June 22 at 11:59 PM.",
        "Dry-run / shakedown: Saturday June 20, 9:00 AM, Clearwater Fire Training Center — open to all members.",
        "Bring your own wired headset (3.5 mm or 1/4 in plug).",
        "If you’ve never operated a contest before — we’ll partner you with an experienced op.",
        "Volunteer for setup (Sat June 27, 8 AM) or teardown (Sun June 28, 2 PM) to learn the gear.",
    ],
)

# 14 - Operator etiquette quick reference
add_bullets(
    "Operating Etiquette — Quick Reference",
    [
        "Don’t forget to “OPON” when starting your shift (TR4W operator change).",
        "Confirm the band/mode is open before calling — listen first.",
        "Log every QSO immediately; the stats screen updates from your log.",
        "If something is wrong with your rig or laptop, retrace the setup card on your table before calling the LAN team.",
        "Have fun — it’s a contest, not an exam.",
    ],
)

# 15 - Q&A
s = prs.slides.add_slide(LAY_SECTION)
s.shapes.title.text = "Q & A"
for r in s.shapes.title.text_frame.paragraphs[0].runs:
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r.font.size = Pt(60)
tb = s.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(2))
tf = tb.text_frame
tf.text = "Questions?  Concerns?  Volunteers?"
for r in tf.paragraphs[0].runs:
    r.font.italic = True
    r.font.color.rgb = MUTED
    r.font.size = Pt(28)
p2 = tf.add_paragraph()
p2.text = "Talk to your section lead: Ryan AF4O (antennas) · Fred W2SUB (education) · Tom NY4I (logging/network) · Paul KC4YDY (VHF) · Ron W4RFA (food)"
for r in p2.runs:
    r.font.color.rgb = MUTED
    r.font.size = Pt(18)

prs.save(OUT)
print(f"Saved {OUT}")
print(f"Slide count: {len(prs.slides)}")
